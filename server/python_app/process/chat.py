from __future__ import annotations
import uuid
import asyncio
import json
import os
import re
import shutil
from typing import Optional

import httpx

from ..config import OPENCODE_BASE_URL
from ..database import execute, execute_write, execute_one, get_pool
from ..agent.registry import registry
from ..agent.base import SendResult

def _stable_hash(s: str, mod: int = 10000) -> int:
    h = 5381
    for c in s:
        h = ((h << 5) + h) + ord(c)
    return abs(h) % mod
from ..control.auth import get_user_api_keys
from ..acp.client import ACPClient
from ..acp.protocol import ACPMethod


_send_locks: dict[str, asyncio.Lock] = {}
_oc_cache: dict[str, str] = {}
_acp_client: Optional[ACPClient] = None
_memory_mgr: Optional["MemoryManager"] = None


def set_acp_client(client: ACPClient):
    global _acp_client
    _acp_client = client


def set_memory_manager(mgr):
    global _memory_mgr
    _memory_mgr = mgr


async def list_sessions(username: str) -> list[dict]:
    rows = await execute(
        """SELECT id, title, model, pinned, created_at, updated_at
           FROM chat_sessions WHERE username=%s
           ORDER BY pinned DESC, updated_at DESC LIMIT 50""",
        (username,),
    )
    return [
        {
            "id": r["id"], "title": r["title"], "model": r["model"],
            "pinned": r["pinned"],
            "created": r["created_at"].strftime("%m-%d %H:%M"),
            "updated": r["updated_at"].isoformat(),
        }
        for r in rows
    ]


async def create_session(username: str, model: str = "") -> dict:
    """创建新会话"""
    if not model:
        for pid, inst in registry.list():
            cfg = registry.get_config(pid)
            if cfg.get("requires_key"):
                continue
            ms = cfg.get("models", [])
            if not ms:
                try:
                    ms = [{"id": m.id} for m in await inst.fetch_models()]
                except Exception:
                    pass
            if ms:
                model = ms[0]["id"]
                break
        if not model:
            for pid, inst in registry.list():
                cfg = registry.get_config(pid)
                ms = cfg.get("models", [])
                if ms:
                    model = ms[0]["id"]
                    break
        if not model:
            model = "deepseek-v4-flash-free"

    sid = f"chat_{uuid.uuid4().hex[:12]}"
    provider_inst = _get_provider_for_model(model)
    oc_id = ""

    if provider_inst and provider_inst.provider_id == "opencode":
        try:
            oc = await provider_inst.create_session()
            oc_id = oc.get("id", "")
            if oc_id:
                await provider_inst.set_model(oc_id, model)
        except Exception:
            pass

    await execute_write(
        "INSERT INTO chat_sessions (id, username, title, model, oc_session_id) VALUES (%s,%s,%s,%s,%s)",
        (sid, username, "新对话", model, oc_id),
    )
    return {"id": sid, "title": "新对话", "model": model, "oc_session_id": oc_id}


async def delete_session(sid: str, username: str):
    await execute_write("DELETE FROM chat_messages WHERE session_id=%s", (sid,))
    await execute_write("DELETE FROM chat_sessions WHERE id=%s AND username=%s", (sid, username))


async def get_messages(sid: str, username: str) -> dict:
    rows = await execute(
        """SELECT id, role, content, thinking, created_at
           FROM chat_messages WHERE session_id=%s ORDER BY id ASC LIMIT 100""",
        (sid,),
    )
    session = await execute_one("SELECT model FROM chat_sessions WHERE id=%s", (sid,))
    return {
        "username": username,
        "model": session["model"] if session else "",
        "messages": [
            {"id": r["id"], "role": r["role"], "content": r["content"],
             "thinking": r.get("thinking"),
             "time": r["created_at"].strftime("%Y/%-m/%-d %H:%M:%S")}
            for r in rows
        ],
    }


async def update_session_title(sid: str, username: str, title: str):
    await execute_write(
        "UPDATE chat_sessions SET title=%s WHERE id=%s AND username=%s",
        (title, sid, username),
    )


async def update_session_model(sid: str, username: str, model: str):
    await execute_write(
        "UPDATE chat_sessions SET model=%s, oc_session_id='' WHERE id=%s AND username=%s",
        (model, sid, username),
    )


async def toggle_pin(sid: str, username: str) -> int:
    pool = await get_pool()
    async with pool.acquire() as conn:
        async with conn.cursor() as cur:
            await cur.execute(
                "SELECT pinned FROM chat_sessions WHERE id=%s AND username=%s FOR UPDATE",
                (sid, username),
            )
            row = await cur.fetchone()
            new_val = None
            if row:
                new_val = 0 if row[0] else 1
                if new_val == 1:
                    await cur.execute(
                        "UPDATE chat_sessions SET pinned=0 WHERE username=%s AND id!=%s",
                        (username, sid),
                    )
                    await cur.execute(
                        "UPDATE chat_sessions SET pinned=1, updated_at=NOW() WHERE id=%s AND username=%s",
                        (sid, username),
                    )
                else:
                    await cur.execute(
                        "UPDATE chat_sessions SET pinned=0 WHERE id=%s AND username=%s",
                        (sid, username),
                    )
            await conn.commit()
    return new_val if row else 0


async def send_message(sid: str, user_id: str, username: str, content: str = "") -> dict:
    """发送消息给 AI，返回 assistant 回复（含问题和选项）。
    
    流程：
    1. 检查 OpenCode 是否有待解决问题 → 用 content 回答
    2. 否则→ 发送 content 作为新 prompt
    3. 进入 agent 循环：自动批准 permission，检测 question（保存为消息返回）或 finish:stop
    4. 始终返回 type: "content"（问题和选项也作为内容返回）
    """
    session = await execute_one(
        "SELECT oc_session_id, model FROM chat_sessions WHERE id=%s", (sid,),
    )
    oc_id = session["oc_session_id"] if session else ""
    session_model = session["model"] if session else ""
    if not session_model:
        return {"success": False, "error": "No model selected"}
    provider_inst = _get_provider_for_model(session_model)
    if not provider_inst:
        return {"success": False, "error": f"Unknown provider for model: {session_model}"}

    lock = _send_locks.setdefault(sid, asyncio.Lock())
    try:
        async with lock:
            # 检查是否有待回答问题 → 用用户消息回答
            current_q = await provider_inst.check_question(oc_id)
            if current_q and current_q.get("questions"):
                qid = current_q.get("id", "")
                await provider_inst.answer_question(oc_id, qid, [[content]])
                print(f"[Process] 回答问题: session={sid[:8]} qid={qid}", flush=True)
            else:
                # 正常发送消息
                if not content:
                    return {"success": False, "error": "No content"}
                await execute_write(
                    "INSERT INTO chat_messages (session_id, role, content) VALUES (%s,%s,%s)",
                    (sid, "user", content),
                )
                count = await execute(
                    "SELECT COUNT(*) as c FROM chat_messages WHERE session_id=%s AND role='user'", (sid,),
                )
                if count and count[0]["c"] == 1:
                    title = content[:30] + ("..." if len(content) > 30 else "")
                    await update_session_title(sid, username, title)
                if content.startswith("<think>"):
                    content = content.replace("<think>", "").replace("</think>", "")
                    content = "\n\n请按以下格式输出：\n<think>\n你的推理过程\n</think>\n你的最终答案\n\n问题：" + content

                # 获取 API Key（如果需要）
                api_key = ""
                sub_key = ""
                cfg = registry.get_config(provider_inst.provider_id)
                if cfg.get("requires_key"):
                    user_keys = await get_user_api_keys(user_id)
                    subs = cfg.get("sub_providers", {})
                    for sk, sv in subs.items():
                        kt = sv.get("key_type", sk)
                        if user_keys.get(kt):
                            api_key = user_keys[kt]; sub_key = sk; break
                    if not api_key:
                        return {"success": False, "error": "API key not configured"}

                if provider_inst.provider_id == "opencode" and not oc_id:
                    oc = await provider_inst.create_session()
                    oc_id = oc.get("id", "")
                    if oc_id:
                        await provider_inst.set_model(oc_id, session_model)
                    await execute_write("UPDATE chat_sessions SET oc_session_id=%s WHERE id=%s", (oc_id, sid))

                print(f"[Process] 发送消息: session={sid[:8]} model={session_model}", flush=True)
                content_with_memory = await _inject_memory_context(user_id, content, sid)
                try:
                    await provider_inst.send_prompt(oc_id, content_with_memory, session_model)
                except RuntimeError as e:
                    return {"success": False, "error": f"Send failed: {e}"}
                print(f"[Process] prompt 发送成功: session={sid[:8]}", flush=True)

            # ====== agent 循环 ======
            print(f"[Process] 进入 agent 循环: session={sid[:8]}", flush=True)
            known_ids = await provider_inst._fetch_known_ids(oc_id)
            result = await _agent_loop(provider_inst, oc_id, sid, known_ids)
            print(f"[Process] agent 循环结束: session={sid[:8]} result_keys={list(result.keys())}", flush=True)

            if result.get("question"):
                # 问题是 assistant 消息的一部分，保存到 DB 后返回
                q_text = _format_questions_as_message(result["questions"])
                await execute_write(
                    "INSERT INTO chat_messages (session_id, role, content) VALUES (%s,%s,%s)",
                    (sid, "assistant", q_text),
                )
                await execute_write("UPDATE chat_sessions SET updated_at=NOW() WHERE id=%s", (sid,))
                return {"success": True, "type": "content", "content": q_text, "thinking": "",
                        "_qid": result["qid"]}

            text = result.get("text", "").strip()
            if text:
                import re as _re
                tm = _re.search(r"<think>(.*?)</think>", text, _re.DOTALL)
                if tm:
                    before = text[:text.index("<think>")].strip()
                    after = text[text.index("</think>") + 8:].strip()
                    tc = tm.group(1).strip()
                    if before and tc.startswith(before): tc = tc[len(before):].strip()
                    tp = (before + "\n\n" + tc).strip()
                    ct = after
                else:
                    tp = ""; ct = text
                await execute_write(
                    "INSERT INTO chat_messages (session_id, role, content, thinking) VALUES (%s,%s,%s,%s)",
                    (sid, "assistant", ct, tp),
                )
                await execute_write("UPDATE chat_sessions SET updated_at=NOW() WHERE id=%s", (sid,))
                if content: await _extract_and_store_memories(user_id, content)
                asyncio.create_task(_export_ppt_files(sid, oc_id, username))
                return {"success": True, "type": "content", "content": ct, "thinking": tp}

            return {"success": False, "error": "AI 返回空响应，请重试"}
    finally:
        stale = [k for k, v in _send_locks.items() if v.locked() and k != sid]
        for k in stale: del _send_locks[k]


async def _agent_loop(provider_inst, oc_id: str, sid: str, known_ids: set[str]) -> dict:
    """Agent 循环：自动批准 permission，检测 question 或 finish:stop 后返回"""
    known_ids = set(known_ids) if not isinstance(known_ids, set) else known_ids
    start = asyncio.get_event_loop().time()
    timeout = 300
    error_count = 0

    async with httpx.AsyncClient(timeout=30) as client:
        while True:
            if asyncio.get_event_loop().time() - start > timeout:
                return {"text": ""}

            try:
                perm = await provider_inst.check_permission(oc_id)
                if perm:
                    print(f"[Agent] 自动批准 permission: {perm['id']} action={perm.get('action','')}")
                    await provider_inst.approve_permission(oc_id, perm["id"])
                    continue
            except Exception:
                pass

            try:
                q = await provider_inst.check_question(oc_id)
                if q and q.get("questions"):
                    print(f"[Agent] 检测到问题: qid={q['id']} questions={len(q['questions'])}")
                    return {"question": True, "qid": q["id"], "questions": q["questions"]}
            except Exception:
                pass

            try:
                resp = await client.get(
                    f"{OPENCODE_BASE_URL}/api/session/{oc_id}/message",
                    params={"from": 0, "to": 100},
                )
                messages = resp.json().get("data", [])
            except Exception:
                await asyncio.sleep(1)
                continue

            for msg in messages:
                if msg.get("id") in known_ids: continue
                if msg.get("type") == "assistant" and msg.get("finish") in ("stop", "error"):
                    text = "".join(b.get("text", "") for b in (msg.get("content") or []) if isinstance(b, dict))
                    if text.strip():
                        return {"text": text.strip()}
                    if msg.get("finish") == "error":
                        error_count += 1
                        if error_count < 3:
                            print(f"[Agent] AI 返回错误，重试 {error_count}/3", flush=True)
                            await asyncio.sleep(3)
                            continue  # 重试
                        return {"text": ""}  # AI error, no text after retries

            await asyncio.sleep(0.5)


def _format_questions_as_message(questions: list) -> str:
    """把 AI 的提问和选项格式化为一条 assistant 消息"""
    parts = []
    for q in questions:
        hdr = q.get("header", "")
        desc = q.get("question", "")
        if hdr: parts.append(f"**{hdr}**")
        if desc: parts.append(desc)
        opts = q.get("options", [])
        for o in opts:
            label = o.get("label", "")
            d = o.get("description", "")
            parts.append(f"- {label}" + (f"：{d}" if d else ""))
        parts.append("")
    return "\n".join(parts).strip()


async def _export_ppt_files(sid: str, oc_id: str, username: str) -> Optional[str]:
    """查找 AI 生成的 PPT 文件，复制到统一路径，渲染 HD 并合成 .pptx"""
    import os, shutil, time
    PPT_BASE = "/tmp/opencode/ai-chat/server/ppt"

    # 找最新创建的 ppt-output 目录
    runs_dir = "/tmp/opencode/ppt-output/runs"
    if not os.path.isdir(runs_dir):
        return None

    candidates = sorted(
        [d for d in os.listdir(runs_dir) if os.path.isdir(os.path.join(runs_dir, d))],
        key=lambda d: os.path.getmtime(os.path.join(runs_dir, d)),
        reverse=True,
    )
    if not candidates:
        return None

    # 取 poll 开始后创建的最新目录（取前 2 个，避免跨 session 误判）
    src = None
    now = time.time()
    for d in candidates:
        dpath = os.path.join(runs_dir, d)
        mtime = os.path.getmtime(dpath)
        # 目录必须在 poll 启动前后 5 分钟内创建
        if abs(now - mtime) < 600:
            src = dpath
            break

    if not src or not os.path.isdir(src):
        return None

    dest_base = os.path.join(PPT_BASE, username, os.path.basename(src))
    os.makedirs(dest_base, exist_ok=True)

    # 复制 slides/ 和 png/（如果有）
    for sub in ("slides", "png"):
        s_sub = os.path.join(src, sub)
        if os.path.isdir(s_sub):
            d_sub = os.path.join(dest_base, sub)
            shutil.rmtree(d_sub, ignore_errors=True)
            shutil.copytree(s_sub, d_sub)

    # 复制支撑文件
    for fn in ("style.json", "outline.txt", "search.txt"):
        fp = os.path.join(src, fn)
        if os.path.isfile(fp):
            shutil.copy2(fp, dest_base)

    # ---- 生成 HD PNG + index.html + .pptx ----
    slides_dir = os.path.join(dest_base, "slides")
    png_dir = os.path.join(dest_base, "png")
    if not os.path.isdir(slides_dir):
        return dest_base

    # 合并 index.html
    basecss = ""
    css_path = os.path.join(slides_dir, "_base.css")
    if os.path.isfile(css_path):
        basecss = open(css_path, encoding="utf-8").read()

    slide_files = sorted(
        [f for f in os.listdir(slides_dir) if f.startswith("slide-") and f.endswith(".html")],
        key=lambda x: int(x.replace("slide-", "").replace(".html", "")),
    )
    if slide_files:
        import re as _re2
        parts = []
        for f in slide_files:
            html = open(os.path.join(slides_dir, f), encoding="utf-8").read()
            m = _re2.search(r"<body>(.*?)</body>", html, _re2.S)
            inner = m.group(1) if m else html
            parts.append(f'<section class="slidewrap">{inner}</section>')

        combined = (
            f'<!DOCTYPE html><html lang="zh-CN"><head><meta charset="utf-8">'
            f'<title>PPT</title><style>{basecss}'
            f'body{{background:#0b0e14;margin:0;padding:20px;font-family:system-ui,"PingFang SC","Microsoft YaHei",sans-serif;}}'
            f'.slidewrap{{margin:0 auto 16px;max-width:1280px;}}'
            f'.slidewrap .slide{{height:720px;}}'
            f'</style></head><body>'
            f'<div style="position:fixed;top:10px;right:16px;z-index:99;background:#11182a;color:#cfe;padding:8px 12px;border-radius:8px;font-size:13px;">'
            f'PPT · 共 {len(slide_files)} 页</div>'
            f'{"".join(parts)}</body></html>'
        )
        with open(os.path.join(dest_base, "index.html"), "w", encoding="utf-8") as f:
            f.write(combined)

    # HD PNG rendering with Playwright
    hd_png_dir = os.path.join(dest_base, "png_hd")
    if slide_files and not os.listdir(hd_png_dir) if os.path.isdir(hd_png_dir) else True:
        try:
            from playwright.sync_api import sync_playwright

            os.makedirs(hd_png_dir, exist_ok=True)
            with sync_playwright() as p:
                browser = p.chromium.launch()
                page = browser.new_page(viewport={"width": 3840, "height": 2160}, device_scale_factor=2)
                for f in slide_files:
                    fp = os.path.join(slides_dir, f)
                    page.goto(f"file://{fp}", wait_until="networkidle")
                    page.wait_for_timeout(300)
                    page.locator(".slide").screenshot(path=os.path.join(hd_png_dir, f.replace(".html", ".png")))
                browser.close()
            print(f"[Export] HD PNG rendered: {len(slide_files)} slides")
        except Exception as e:
            print(f"[Export] HD PNG failed: {e}")
            # fallback: copy existing png
            if os.path.isdir(png_dir):
                for f in os.listdir(png_dir):
                    shutil.copy2(os.path.join(png_dir, f), os.path.join(hd_png_dir, f))

    # Generate .pptx
    hd_src = hd_png_dir if os.path.isdir(hd_png_dir) and os.listdir(hd_png_dir) else png_dir
    if os.path.isdir(hd_src):
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.dml.color import RGBColor

            png_files = sorted(
                [f for f in os.listdir(hd_src) if f.endswith(".png")],
                key=lambda x: int(x.replace("slide-", "").replace(".png", "")),
            )
            if png_files:
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                for fn in png_files:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(0x0b, 0x0e, 0x14)
                    slide.shapes.add_picture(os.path.join(hd_src, fn), Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                pptx_path = os.path.join(dest_base, "presentation-hd.pptx")
                prs.save(pptx_path)
                print(f"[Export] PPTX saved: {pptx_path} ({os.path.getsize(pptx_path)//1024}KB)")
        except Exception as e:
            print(f"[Export] PPTX failed: {e}")

    return dest_base


_MEMORY_CATEGORIES = {
    "preference": {"label": "偏好", "ttl": 86400 * 365 * 10},  # 10年 ≈ 永久
    "fact": {"label": "事实", "ttl": 86400 * 365 * 10},
    "task": {"label": "任务", "ttl": 3600},  # 1小时 ≈ 当前会话
    "rule": {"label": "规则", "ttl": 86400 * 365 * 10},
}


async def _inject_memory_context(user_id: str, content: str, sid: str) -> str:
    """从记忆库召回用户信息 + 全局信息，注入到当前消息前导中"""
    if not _memory_mgr:
        return content

    memories = []
    for scope in ("long_term", "short_term"):
        items = await _memory_mgr.list_user_memories(user_id, scope)
        for item in items:
            if item["key"].startswith("_") or item["key"] == "injected":
                continue
            # key 格式: user:{user_id}:{cat}:{subcat}[:{session_id}]:{hash}
            parts = item["key"].split(":")
            cat = parts[2] if len(parts) > 2 else "general"
            subcat = parts[3] if len(parts) > 3 else ""
            # 任务按 session 隔离：key 中必须有当前 sid 才注入
            if cat == "task":
                if len(parts) < 6 or parts[4] != sid:
                    continue
            label = _MEMORY_CATEGORIES.get(cat, {}).get("label", cat)
            display_key = f"{subcat}" if subcat else cat
            if cat == "rule":
                memories.append(f"[{label}] {item['value']}")
            else:
                memories.append(f"[{label}] {display_key}: {item['value']}")

    global_items = await _memory_mgr.list_global()
    for item in global_items:
        memories.append(f"[全局] {item['key']}: {item['value']}")

    if not memories:
        return content

    note = "\n".join(memories[:8])
    tagged = f"[Memory Note: 以下是与用户相关的已知信息，供参考]\n{note}\n\n{content}"
    return tagged[:8000]


async def _extract_and_store_memories(user_id: str, content: str, session_id: str = ""):
    """从用户消息中提取有价值的记忆并存储"""
    if not _memory_mgr or len(content) > 500:
        return

    patterns = [
        (r"(?:全局记住|全局规则)[：:]?\s*(.{5,200})", "global:rule"),
        (r"(?:我是|我叫|我(?:是|做)(?:一个|名)?)(.{2,20}(?:工程师|开发|设计|产品|运营|学生|教师))", "fact:role"),
        (r"(?:我叫|我的名字是)([^什哪谁什么名字]{1,10})", "fact:name"),
        (r"(?:我喜欢|我习惯|我偏好|我更(?:喜欢|倾向))(.{2,30})", "preference:style"),
        (r"(?:我在|我?工作在|我任职于)(.{2,30}(?:公司|单位|团队))", "fact:company"),
        (r"(?:我用|我(?:正在)?学|我熟悉|我擅长|我喜欢用|我经常用|我习惯用)(.{2,30}(?:语言|框架|工具|技术|系统|平台))", "preference:tech"),
        (r"(?:我喜欢用|我用|我经常用|我习惯用)(.{2,30})", "preference:tech"),
        (r"(?:请记住|记住|以后记住)[：:]?\s*(.{5,200})", "rule:instruction"),
    ]

    for pattern, key in patterns:
        m = re.search(pattern, content)
        if m:
            val = m.group(1).strip()[:100]
            if key.startswith("global:"):
                try:
                    await _memory_mgr.set_global(key.split(":", 1)[1], val)
                except Exception:
                    pass
                continue
            cat = key.split(":")[0]
            ttl = _MEMORY_CATEGORIES.get(cat, {}).get("ttl", 3600)
            # 任务 key 嵌入 session_id 实现会话隔离
            if cat == "task" and session_id:
                store_key = f"user:{user_id}:{key}:{session_id}:{_stable_hash(val)}"
            else:
                store_key = f"user:{user_id}:{key}:{_stable_hash(val)}"
            await _memory_mgr.remember(store_key, val, scope="short_term", ttl=3600)
            if cat != "task":
                await _memory_mgr.remember(store_key, val, scope="long_term", ttl=ttl)
            break


def _get_provider_for_model(model_id: str):
    """根据模型 ID 找到对应的 provider 实例"""
    # 检查已知模型的前缀映射
    for pid, inst in registry.list():
        if model_id in [m["id"] for m in registry.get_config(pid).get("models", []) if isinstance(m, dict)]:
            return inst
        # 尝试通过 provider 自己获取所有模型来匹配
        try:
            cfg = registry.get_config(pid)
            for m in cfg.get("models", []):
                if m["id"] == model_id:
                    return inst
        except Exception:
            pass
    # fallback: opencode 兜底
    return registry.get("opencode")
